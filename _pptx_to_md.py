"""
One-off: extract masterclass.pptx to masterclass.md with images in assets/.
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape

ROOT = Path(__file__).resolve().parent
PPTX = ROOT / "masterclass.pptx"
OUT_MD = ROOT / "masterclass.md"
ASSETS = ROOT / "assets"


def _merge_runs_text(runs) -> str:
    chunks: list[str] = []
    for run in runs:
        t = run.text
        if not t:
            continue
        if chunks and not chunks[-1].endswith(("\n", " ", "\t")) and not t[:1].isspace():
            a = chunks[-1][-1:] if chunks[-1] else ""
            b = t[:1]
            if a.isalnum() and b.isalnum():
                chunks.append(" ")
        chunks.append(t)
    return "".join(chunks).strip()


def _text_from_shape(shape) -> str:
    parts: list[str] = []
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            line = _merge_runs_text(p.runs)
            if line:
                parts.append(line)
    return "\n\n".join(parts).strip()


def _iter_shapes(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        gs: GroupShape = shape  # type: ignore[assignment]
        for s in gs.shapes:
            yield from _iter_shapes(s)
    else:
        yield shape


def _slide_blocks(slide, si: int, img_counter: list[int]) -> list[str]:
    blocks: list[str] = []
    for shape in slide.shapes:
        for sub in _iter_shapes(shape):
            if sub.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_counter[0] += 1
                n = img_counter[0]
                ext = sub.image.ext  # type: ignore[attr-defined]
                if not ext:
                    ext = "png"
                name = f"pptx-slide-{si:02d}-img-{n:02d}.{ext}"
                path = ASSETS / name
                path.write_bytes(sub.image.blob)  # type: ignore[attr-defined]
                blocks.append(f"![](assets/{name})")
            else:
                t = _text_from_shape(sub)
                if t:
                    blocks.append(t)
    # De-duplicate consecutive duplicates only
    out: list[str] = []
    prev = None
    for b in blocks:
        if b == prev:
            continue
        prev = b
        out.append(b)
    return out


def main() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(PPTX))
    lines: list[str] = []
    lines.append("# SGR Agent Core — Мастер-класс")
    lines.append("")
    lines.append("_Автоматическая конвертация из PPTX в Markdown._")
    lines.append("")
    lines.append(f"- Исходный файл: `{PPTX.name}`")
    lines.append(f"- Слайдов: {len(prs.slides)}")
    lines.append("")
    lines.append("---")
    lines.append("")

    img_counter = [0]
    for si, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Слайд {si}")
        lines.append("")
        blocks = _slide_blocks(slide, si, img_counter)
        for b in blocks:
            lines.append(b)
            lines.append("")
        lines.append("---")
        lines.append("")

    text = "\n".join(lines).rstrip() + "\n"
    text = re.sub(r"\n{3,}", "\n\n", text)
    OUT_MD.write_text(text, encoding="utf-8")
    print(f"Wrote {OUT_MD} ({len(prs.slides)} slides, {img_counter[0]} images)")


if __name__ == "__main__":
    main()
